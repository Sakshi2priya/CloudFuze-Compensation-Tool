"""
Commission policy presentation: structured team cards (from ``commission_policy`` / JSON)
and optional PDF/PPTX text extraction for the Rules & Policy page.
"""

from __future__ import annotations

import base64
import html
import io
from pathlib import Path
from typing import Any

import streamlit as st

import commission_policy as cp


TEAM_SKIN = {
    "SMB": {"border": "#1565c0", "tint1": "#e3f2fd", "chip": "#1565c0"},
    "AM": {"border": "#00897b", "tint1": "#e0f2f1", "chip": "#00695c"},
    "ENT": {"border": "#7b1fa2", "tint1": "#f3e5f5", "chip": "#6a1b9a"},
    "OB": {"border": "#f57c00", "tint1": "#fff8e1", "chip": "#ef6c00"},
    "MN": {"border": "#0d9488", "tint1": "#ccfbf1", "chip": "#0f766e"},
}


def _esc(s: Any) -> str:
    return html.escape(str(s if s is not None else ""))


def _display_name_from_token(tok: str) -> str:
    """Policy matching tokens → readable names (sentence case / dotted emails)."""
    raw = (tok or "").strip()
    if not raw:
        return ""
    if "." in raw and "@" not in raw:
        return ".".join(_display_name_from_token(p) for p in raw.split(".") if p)
    if "_" in raw:
        return " ".join(_display_name_from_token(p) for p in raw.split("_") if p)
    return raw[0].upper() + raw[1:].lower() if raw else ""


def _dedupe_tokens_case_insensitive(tokens: list[Any]) -> list[str]:
    out: list[str] = []
    seen: set[str] = set()
    for t in tokens:
        s = str(t).strip()
        if not s:
            continue
        k = s.lower()
        if k in seen:
            continue
        seen.add(k)
        out.append(s)
    return out


# Snapshot tokens that are alternate spellings → single display name (PPT-style).
_SNAPSHOT_REP_DISPLAY_CANON: dict[str, str] = {
    "krithika": "Kritika",
    "karthik": "Kartik",
}


def _snapshot_rep_display_list(tokens: list[str]) -> list[str]:
    """Order-preserving list of rep display names; merges known alternate spellings."""
    out: list[str] = []
    seen: set[str] = set()
    for t in _dedupe_tokens_case_insensitive(tokens):
        low = str(t).strip().lower()
        disp = _SNAPSHOT_REP_DISPLAY_CANON.get(low) or _display_name_from_token(str(t))
        key = disp.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(disp)
    return out


def _fmt_money(x: float) -> str:
    try:
        return f"${float(x):,.0f}"
    except (TypeError, ValueError):
        return str(x)


def extract_policy_document_text(file_bytes: bytes, filename: str) -> tuple[str | None, str | None]:
    """Return (plain_text, error_message). Supports PDF and PPTX."""
    name = (filename or "").lower()
    if name.endswith(".pdf"):
        try:
            from pypdf import PdfReader
        except ImportError:
            return None, "Missing dependency: install **pypdf**."
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            chunks: list[str] = []
            for page in reader.pages:
                chunks.append(page.extract_text() or "")
            return "\n\n".join(chunks), None
        except Exception as e:
            return None, str(e)
    if name.endswith(".pptx"):
        try:
            from pptx import Presentation
        except ImportError:
            return None, "Missing dependency: install **python-pptx** to read PowerPoint files."
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            slides_out: list[str] = []
            for slide in prs.slides:
                parts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = (getattr(shape, "text", None) or "").strip()
                        if t:
                            parts.append(t)
                slides_out.append("\n".join(parts))
            return "\n\n— slide —\n\n".join(slides_out), None
        except Exception as e:
            return None, str(e)
    return None, "Please upload a **.pdf** or **.pptx** file."


def bucket_extracted_lines(text: str) -> dict[str, list[str]]:
    """Route extracted lines into SMB / AM / ENT / Outbound / Other (heuristic)."""
    lines = [ln.strip() for ln in text.replace("\r", "\n").split("\n")]
    lines = [ln for ln in lines if len(ln) > 1]
    buckets = {k: [] for k in ("SMB", "AM", "ENT", "OB", "OTHER")}
    for ln in lines:
        low = ln.lower()
        ob_hit = any(
            p in low
            for p in (
                "outbound",
                "meeting incentive",
                "booked in nam",
                "western europe",
                "north america",
                "outbound meeting",
            )
        )
        if ob_hit and "account management" not in low:
            buckets["OB"].append(ln)
        elif "enterprise" in low or "anthony" in low:
            buckets["ENT"].append(ln)
        elif any(
            p in low
            for p in (
                "account management",
                " vivin",
                " joy ",
                "arundhat",
            )
        ):
            buckets["AM"].append(ln)
        elif any(p in low for p in (" smb", "smb ", "smb:", "chitradip", "group a", "group b", " smb_")):
            buckets["SMB"].append(ln)
        elif low.startswith("smb"):
            buckets["SMB"].append(ln)
        else:
            buckets["OTHER"].append(ln)
    return buckets


def _tier_pct_range(t: dict[str, Any]) -> str:
    try:
        mn = float(t.get("min_pct")) if t.get("min_pct") is not None else None
    except (TypeError, ValueError):
        mn = None
    try:
        mx = float(t.get("max_pct")) if t.get("max_pct") is not None else None
    except (TypeError, ValueError):
        mx = None
    if mn is None:
        return "—"
    if mx is None:
        return f"{mn:g}%+"
    return f"{mn:g}%–{mx:g}%"


def _policy_hub_header_html() -> str:
    """Prominent policy name + subtitle (user-requested ‘name’ block)."""
    pn = _esc(cp.POLICY_NAME)
    return f"""<div style="margin:4px 0 8px 0;padding:22px 26px;border-radius:18px;background:linear-gradient(118deg,#e8f4fc 0%,#f3e8ff 42%,#fff8e6 100%);border:1px solid #c7d2fe;box-shadow:0 6px 24px rgba(15,23,42,.07);">
<div style="font-size:.72rem;font-weight:800;letter-spacing:.14em;color:#64748b;text-transform:uppercase;">Policy name</div>
<div style="font-size:clamp(1.35rem,2.4vw,1.75rem);font-weight:800;color:#0f172a;margin:8px 0 10px 0;line-height:1.25;">{pn}</div>
<div style="font-size:.96rem;color:#475569;line-height:1.5;">Rules &amp; policy hub</div>
</div>"""


def _policy_hub_nav_css() -> str:
    """Extra vertical rhythm around the policy section navigator (segmented control)."""
    return """<style>
    /* Policy hub: breathing room around segmented control + taller segments */
    div[data-testid="stMainBlockContainer"] div[data-baseweb="tab-list"],
    div[data-testid="stMainBlockContainer"] [data-testid="stHorizontalBlock"]:has(button[kind="header"]) {
        gap: 12px;
    }
    /* Segmented control (Streamlit ≥1.41): stretch + padding inside buttons */
    div[data-testid="stMainBlockContainer"] [data-baseweb="segmented-control"] {
        margin-top: 6px !important;
        margin-bottom: 26px !important;
        min-height: 52px !important;
    }
    div[data-testid="stMainBlockContainer"] [data-baseweb="segmented-control"] button {
        padding: 12px 18px !important;
        min-height: 48px !important;
        font-size: 0.95rem !important;
        line-height: 1.35 !important;
    }
    </style>"""


def _team_wrap(team_key: str, ribbon: str, headline: str, inner_html: str) -> str:
    s = TEAM_SKIN[team_key]
    return f"""<div style="margin-bottom:34px;border-radius:16px;overflow:hidden;box-shadow:0 6px 22px rgba(15,23,42,.1);
border-left:6px solid {s["border"]};background:linear-gradient(100deg,{s["tint1"]} 0%,#ffffff 52%);">
<div style="padding:16px 20px 4px 20px;">
<span style="display:inline-block;padding:5px 12px;border-radius:999px;background:{s["chip"]};color:#fff;font-size:.72rem;font-weight:800;letter-spacing:.08em;">{_esc(ribbon)}</span>
<h3 style="margin:12px 0 0 0;color:#0f172a;font-size:1.28rem;font-weight:800;">{_esc(headline)}</h3>
</div>
<div style="padding:8px 20px 20px 20px;color:#334155;font-size:.94rem;line-height:1.6;">{inner_html}</div>
</div>"""


def _snapshot_group_display_line(group: dict[str, Any], fallback_tokens: list[str]) -> str:
    """Prefer ``policy_display_names`` from snapshot JSON; else derive from ``name_tokens``."""
    raw = group.get("policy_display_names")
    if isinstance(raw, list):
        names = [str(x).strip() for x in raw if str(x).strip()]
        if names:
            return ", ".join(names)
    return ", ".join(_snapshot_rep_display_list(fallback_tokens))


def _smb_group_rep_names_html() -> str:
    """Group A / B rep lists from quota snapshot (title case), when snapshot mode is on."""
    if not cp.SMB_USE_QUOTA_SNAPSHOT or not cp.SMB_QUOTA_SNAPSHOT_RAW:
        return ""
    snap = cp.SMB_QUOTA_SNAPSHOT_RAW
    if not isinstance(snap, dict):
        return ""
    ga = snap.get("group_a") if isinstance(snap.get("group_a"), dict) else {}
    gb = snap.get("group_b") if isinstance(snap.get("group_b"), dict) else {}
    a_toks = list(ga.get("name_tokens") or [])
    b_toks = list(gb.get("name_tokens") or [])
    a_names = _snapshot_group_display_line(ga, a_toks)
    b_names = _snapshot_group_display_line(gb, b_toks)
    qa = ga.get("quota_usd_each")
    qb = gb.get("quota_usd_each")
    q_note = ""
    if qa is not None and qb is not None:
        q_note = (
            f"<p style='margin:0 0 10px 0;font-size:.88rem;color:#64748b;'>Per-rep quota in snapshot: "
            f"Group A {_fmt_money(float(qa))} each · Group B {_fmt_money(float(qb))} each.</p>"
        )
    elif qa is not None:
        q_note = f"<p style='margin:0 0 10px 0;font-size:.88rem;color:#64748b;'>Group A quota each: {_fmt_money(float(qa))}.</p>"
    elif qb is not None:
        q_note = f"<p style='margin:0 0 10px 0;font-size:.88rem;color:#64748b;'>Group B quota each: {_fmt_money(float(qb))}.</p>"
    return (
        "<div style='margin:16px 0 20px 0;padding:16px 18px;border-radius:14px;background:#f8fafc;border:1px solid #e2e8f0;'>"
        "<p style='margin:0 0 12px 0;font-weight:800;color:#0f172a;font-size:1.02rem;'>SMB Group A &amp; Group B — reps</p>"
        f"{q_note}"
        "<p style='margin:0 0 8px 0;'><strong style='color:#1565c0;'>Group A</strong></p>"
        f"<p style='margin:0 0 14px 0;line-height:1.65;font-size:1.02rem;'>{_esc(a_names) if a_names else '—'}</p>"
        "<p style='margin:0 0 8px 0;'><strong style='color:#1565c0;'>Group B</strong></p>"
        f"<p style='margin:0;line-height:1.65;font-size:1.02rem;'>{_esc(b_names) if b_names else '—'}</p>"
        "</div>"
    )


def _html_table(headers: list[str], rows: list[list[str]], head_bg: str) -> str:
    th = "".join(f"<th style='text-align:left;padding:8px 10px;border-bottom:1px solid #e2e8f0;'>{_esc(h)}</th>" for h in headers)
    body = []
    for r in rows:
        tds = "".join(f"<td style='padding:7px 10px;border-bottom:1px solid #f1f5f9;'>{c}</td>" for c in r)
        body.append(f"<tr>{tds}</tr>")
    return f"<table style='width:100%;border-collapse:collapse;font-size:.9rem;margin-top:6px;'><thead><tr style='background:{head_bg};'>{th}</tr></thead><tbody>{''.join(body)}</tbody></table>"


def _mandatory_commission_rules_html(min_pct: float, *, microsoft_google_license: bool = False) -> str:
    """Mandatory bullets for SMB / AM (wording aligned with policy deck)."""
    b1 = (
        f"A minimum of <strong>{min_pct:.0f}%</strong> achievement of the individual quota is mandatory to qualify "
        "for commission payout."
    )
    items = [f"<li style='margin:8px 0;line-height:1.5;'>{b1}</li>"]
    if microsoft_google_license:
        items.append(
            "<li style='margin:8px 0;line-height:1.5;'><strong>No commission:</strong> will be paid on the reselling of "
            "Microsoft or Google licenses.</li>"
        )
    return (
        "<p style='margin-top:18px;margin-bottom:8px;'><strong>Mandatory</strong></p>"
        f"<ul style='margin:0 0 4px 0;padding-left:1.35rem;list-style-type:disc;'>{''.join(items)}</ul>"
    )


def build_smb_inner_html() -> str:
    parts: list[str] = []
    tgt = float(cp.TEAM_QUARTERLY_TARGETS_USD.get(cp.SMB_TEAM_NAME, 0) or 0)
    parts.append(f"<p style='margin-top:0;'><strong>Quarterly team revenue target:</strong> {_fmt_money(tgt)}</p>")
    parts.append(_smb_group_rep_names_html())
    if cp.SMB_QUOTA_ACHIEVEMENT_ENABLED:
        parts.append(
            "<p style='margin-top:4px;'><strong>Rep incentives:</strong> Commission % depends on "
            "<strong>individual quota achievement</strong> (Group A vs Group B tiers in the table below).</p>"
        )
        if cp.SMB_INDIVIDUAL_QUOTAS_USD:
            rows = [
                [_esc(k), f"<strong>{_fmt_money(float(v))}</strong>"]
                for k, v in sorted(cp.SMB_INDIVIDUAL_QUOTAS_USD.items(), key=lambda x: x[0])
            ]
            parts.append(_html_table(["Compensation group", "Individual quota (USD)"], rows, "#dbeafe"))
        if cp.SMB_ACHIEVEMENT_TIERS:
            rws = []
            for t in cp.SMB_ACHIEVEMENT_TIERS:
                rws.append(
                    [
                        _tier_pct_range(t),
                        f"{_esc(t.get('group_a_pct'))}%",
                        f"{_esc(t.get('group_b_pct'))}%",
                    ]
                )
            parts.append(
                "<p style='margin-bottom:4px;'><strong>Achievement tiers (rep)</strong></p>"
                + _html_table(["Achievement vs quota", "Group A commission", "Group B commission"], rws, "#bfdbfe")
            )
        if cp.SMB_MANAGER_CHITRADIP_TEAM_TIERS:
            parts.append(
                "<p style='margin-top:20px;margin-bottom:8px;line-height:1.55;'><strong>SMB manager (Chitradip):</strong> "
                "Team revenue vs team goal — commission % by team achievement band.</p>"
            )
            mgr = [[_tier_pct_range(t), f"<strong>{_esc(t.get('commission_pct'))}%</strong>"] for t in cp.SMB_MANAGER_CHITRADIP_TEAM_TIERS]
            parts.append(_html_table(["Team achievement", "Manager commission %"], mgr, "#e0f2fe"))
        parts.append(_mandatory_commission_rules_html(float(cp.SMB_MIN_ACHIEVEMENT_FOR_COMMISSION_PCT)))
    else:
        parts.append("<p><strong>Rep model:</strong> Revenue commission slabs (and optional SMB_A / SMB_B slab sets in JSON).</p>")
        rws = []
        for lo, hi, pct in cp.REP_SLAB_ROWS_FOR_DB:
            band = f"{_fmt_money(lo)} – {_fmt_money(hi) if hi is not None else 'and above'}"
            rws.append([band, f"<strong>{pct:g}%</strong>"])
        parts.append(_html_table(["Closed revenue band (USD)", "Commission %"], rws, "#dbeafe"))
        parts.append(_mandatory_commission_rules_html(float(cp.MIN_INDIVIDUAL_QUOTA_ACHIEVEMENT_PCT)))
    return "".join(parts)


def build_am_inner_html() -> str:
    parts: list[str] = []
    tgt = float(cp.TEAM_QUARTERLY_TARGETS_USD.get(cp.ACCOUNT_MANAGEMENT_TEAM_NAME, 0) or 0)
    parts.append(f"<p style='margin-top:0;'><strong>Quarterly team revenue target:</strong> {_fmt_money(tgt)}</p>")
    if not cp.AM_QUOTA_ACHIEVEMENT_ENABLED:
        parts.append("<p>Account Management quota-achievement mode is <strong>disabled</strong> in JSON.</p>")
        return "".join(parts)
    iq_rows = []
    for row in cp.AM_INDIVIDUAL_QUOTA_ROWS:
        q = row.get("quota_usd")
        toks = _dedupe_tokens_case_insensitive(list(row.get("name_tokens") or []))
        disp = ", ".join(_display_name_from_token(t) for t in toks)
        iq_rows.append([disp if disp else "—", f"<strong>{_fmt_money(float(q))}</strong>" if q is not None else "—"])
    if iq_rows:
        parts.append(_html_table(["Rep (name match)", "Individual quota (USD)"], iq_rows, "#ccfbf1"))
    if cp.AM_ACHIEVEMENT_TIERS:
        rws = []
        for t in cp.AM_ACHIEVEMENT_TIERS:
            rws.append(
                [
                    _tier_pct_range(t),
                    f"{_esc(t.get('team_pct'))}%",
                    f"{_esc(t.get('joy_pct'))}%",
                ]
            )
        parts.append(
            "<p style='margin-bottom:4px;'><strong>Achievement tiers</strong></p>"
            + _html_table(["Achievement vs quota", "Team commission %", "Joy commission %"], rws, "#99f6e4")
        )
    parts.append(
        _mandatory_commission_rules_html(
            float(cp.AM_MIN_ACHIEVEMENT_FOR_COMMISSION_PCT),
            microsoft_google_license=True,
        )
    )
    return "".join(parts)


def build_ent_inner_html() -> str:
    parts: list[str] = []
    tgt = float(cp.TEAM_QUARTERLY_TARGETS_USD.get(cp.ENTERPRISE_TEAM_NAME, 0) or 0)
    parts.append(f"<p style='margin-top:0;'><strong>Quarterly team revenue target (Enterprise):</strong> {_fmt_money(tgt)}</p>")
    if not cp.ENTERPRISE_QUOTA_ACHIEVEMENT_ENABLED:
        parts.append("<p>Enterprise named-rep quota tiers are <strong>disabled</strong> in JSON; default slabs may apply.</p>")
        return "".join(parts)
    for rep_block in cp.ENT_ENTERPRISE_REP_TIERS:
        toks = _dedupe_tokens_case_insensitive(list(rep_block.get("name_tokens") or []))
        disp = ", ".join(_display_name_from_token(t) for t in toks)
        parts.append(f"<p><strong>Named reps:</strong> {_esc(disp)}</p>")
        mna = rep_block.get("min_achievement_pct_for_commission")
        if mna is not None:
            parts.append(f"<p>Minimum achievement for commission: <strong>{_esc(mna)}%</strong></p>")
        tiers = rep_block.get("achievement_tiers") or []
        rws = []
        for t in tiers:
            if isinstance(t, dict):
                rws.append([_tier_pct_range(t), f"<strong>{_esc(t.get('commission_pct'))}%</strong>"])
        if rws:
            parts.append(_html_table(["Achievement vs individual quota", "Commission %"], rws, "#f3e8ff"))
    return "".join(parts)


def build_outbound_inner_html() -> str:
    """Outbound Team: eligibility, geography, who qualifies, how payouts are recorded (matches Admin → Outbound)."""
    parts: list[str] = []
    parts.append(
        f"<p style='margin-top:0;'><strong>{_esc(cp.OUTBOUND_POLICY_LABEL)}</strong> — outbound meeting incentives.</p>"
    )

    payout_rows = [
        [_esc(r.get("meetings", "")), _esc(r.get("payout", ""))] for r in (cp.OUTBOUND_MEETING_PAYOUT_ROWS or [])
    ]
    if payout_rows:
        parts.append(
            f"<p style='margin:18px 0 8px 0;font-weight:800;color:#c2410c;'>{_esc(cp.OUTBOUND_MEETING_PAYOUT_TITLE)}</p>"
            + _html_table(["Meetings (count)", "Payout"], payout_rows, "#ffedd5")
        )
        if (cp.OUTBOUND_MEETING_PAYOUT_NOTE or "").strip():
            parts.append(
                f"<p style='font-size:.9rem;color:#57534e;margin:10px 0 4px 0;line-height:1.5;'>{_esc(cp.OUTBOUND_MEETING_PAYOUT_NOTE)}</p>"
            )

    return "".join(parts)


def build_monthly_plan_inner_html() -> str:
    """Monthly compensation plan — SMB & Account Management. Source: 2026 Monthly Plan deck."""
    import json as _json
    from pathlib import Path as _Path

    parts: list[str] = []
    parts.append(
        "<p style='margin-top:0;'><strong>Sales Commission Policy 2026 — Monthly Plan</strong> "
        "(applies to SMB and Account Management). Quarterly plan above remains in force for quarter-level reporting; "
        "this section governs month-level payouts.</p>"
    )

    # Load both pinned JSON files for the April month example.
    smb_path = _Path(__file__).parent / "policy" / "smb_april_2026_fixed.json"
    am_path = _Path(__file__).parent / "policy" / "am_april_2026_fixed.json"
    smb_data: dict = {}
    am_data: dict = {}
    try:
        if smb_path.exists():
            smb_data = _json.loads(smb_path.read_text(encoding="utf-8"))
    except Exception:
        smb_data = {}
    try:
        if am_path.exists():
            am_data = _json.loads(am_path.read_text(encoding="utf-8"))
    except Exception:
        am_data = {}

    # --- SMB Monthly Quotas table ---
    smb_reps = smb_data.get("reps") or []
    if smb_reps:
        rows = [[_esc(r.get("name", "")), f"<strong>{_fmt_money(float(r.get('target_usd') or 0))}</strong>"] for r in smb_reps]
        smb_total = sum(float(r.get("target_usd") or 0) for r in smb_reps)
        rows.append(["<strong>Total</strong>", f"<strong>{_fmt_money(smb_total)}</strong>"])
        parts.append(
            "<p style='margin:14px 0 6px 0;font-weight:700;color:#0f766e;'>SMB Team — Monthly Quotas</p>"
            + _html_table(["Rep", "April 2026"], rows, "#ccfbf1")
        )

    # --- AM Monthly Quotas table ---
    am_reps = am_data.get("reps") or []
    if am_reps:
        rows = [[_esc(r.get("name", "")), f"<strong>{_fmt_money(float(r.get('target_usd') or 0))}</strong>"] for r in am_reps]
        am_total = sum(float(r.get("target_usd") or 0) for r in am_reps)
        rows.append(["<strong>Total</strong>", f"<strong>{_fmt_money(am_total)}</strong>"])
        parts.append(
            "<p style='margin:14px 0 6px 0;font-weight:700;color:#0f766e;'>Account Management Team — Monthly Quotas</p>"
            + _html_table(["Rep", "April 2026"], rows, "#ccfbf1")
        )

    # --- Individual rep commission tiers ---
    tiers = (smb_data.get("monthly_tiers") or am_data.get("monthly_tiers") or {})
    ind_tiers = tiers.get("individual_tiers") or []
    if ind_tiers:
        ind_rows = []
        for t in ind_tiers:
            lo = t.get("min_pct")
            hi = t.get("max_pct")
            if hi is None:
                band = f"<strong>{lo}%+</strong>"
            elif lo == 0:
                # Per deck wording, the lowest band displays as "<59.9%" (i.e. anyone below
                # the 60% qualifying threshold). Functional threshold stays at 60% via min_pct
                # of the next tier.
                band = "<strong>&lt; 59.9%</strong>"
            else:
                band = f"<strong>{lo}–{hi}%</strong>"
            ind_rows.append([band, f"<strong>{t.get('commission_pct', 0)}%</strong>"])
        parts.append(
            "<p style='margin:18px 0 6px 0;font-weight:700;color:#0f766e;'>Individual rep commission</p>"
            + _html_table(["% of Quota", "Commission %"], ind_rows, "#ccfbf1")
        )

    # --- Manager commission tiers ---
    mgr_tiers = tiers.get("manager_tiers") or []
    if mgr_tiers:
        mgr_rows = []
        for t in mgr_tiers:
            lo = t.get("min_pct")
            hi = t.get("max_pct")
            band = f"<strong>{lo}–{hi}%</strong>" if hi is not None else f"<strong>{lo}%+</strong>"
            mgr_rows.append([band, f"<strong>{t.get('commission_pct', 0)}%</strong>"])
        parts.append(
            "<p style='margin:18px 0 6px 0;font-weight:700;color:#0f766e;'>Manager commission</p>"
            + _html_table(["Team Achievement %", "Commission %"], mgr_rows, "#ccfbf1")
        )

    # --- Special rules ---
    min_ach = tiers.get("min_achievement_pct_for_commission") or 60
    mgr_min = tiers.get("manager_team_minimum_pct") or 60
    md_pct = tiers.get("manage_deal_pct_am_only") or 5
    parts.append(
        "<p style='margin:18px 0 6px 0;font-weight:700;color:#0f766e;'>Special rules</p>"
        "<ul style='margin:4px 0 0 18px;padding:0;line-height:1.65;'>"
        f"<li>Individual reps below <strong>59.9%</strong> of quota earn <strong>0%</strong> commission.</li>"
        f"<li><strong>Managed-deal incentive:</strong> {md_pct}% on managed deals — <strong>AM team only</strong>.</li>"
        f"<li><strong>Manager commission</strong> applies only if the team achieves at least <strong>{int(mgr_min)}%</strong> of the monthly target (no manager commission below {int(mgr_min)}%).</li>"
        "</ul>"
    )

    parts.append(
        "<p style='font-size:.9rem;color:#57534e;margin:14px 0 0 0;'>"
        "Pinned monthly numbers and achievements live in "
        "<code>policy/smb_&lt;month&gt;_&lt;year&gt;_fixed.json</code> and "
        "<code>policy/am_&lt;month&gt;_&lt;year&gt;_fixed.json</code>. Edit those files (or wire a HubSpot fetch) to populate per-rep achievements."
        "</p>"
    )

    return "".join(parts)


def build_shared_manager_pool_html() -> str:
    """Team manager incentive: achievement vs team goal → commission % (global tiers)."""
    rws = []
    for threshold, pct in cp.TEAM_ACHIEVEMENT_COMMISSION_THRESHOLDS_PCT:
        rws.append([f"Team revenue ÷ goal <strong>&gt; {threshold:g}%</strong>", f"<strong>{pct:g}%</strong>"])
    intro = (
        "<p style='margin-top:0;'>Used for <strong>manager / team pool</strong> incentives when team revenue beats "
        "thresholds against the quarterly team goal (see calculator / team incentives).</p>"
    )
    return intro + _html_table(["Condition (team achievement)", "Commission %"], rws, "#e2e8f0")


def render_extracted_bucket_column(team_key: str, title: str, lines: list[str]) -> None:
    skin = TEAM_SKIN[team_key]
    if not lines:
        st.caption(f"*{title} — no matching lines in this extract.*")
        return
    shown = lines[:80]
    more = len(lines) - len(shown)
    body = "<br/>".join(_esc(x) for x in shown)
    if more > 0:
        body += f"<br/><span style='opacity:.75;font-size:.85rem;'>+ {more} more lines…</span>"
    st.markdown(
        f"""<div style="border-radius:14px;border:1px solid {skin["border"]};background:{skin["tint1"]};
        padding:16px 18px;max-height:440px;overflow-y:auto;font-size:.9rem;line-height:1.5;color:#1e293b;">
        <div style="font-weight:800;color:{skin["chip"]};margin-bottom:12px;font-size:1rem;">{_esc(title)}</div>{body}</div>""",
        unsafe_allow_html=True,
    )


def render_commission_policy_page(pdf_path: Path | None, key_prefix: str) -> None:
    """Main Rules & Policy hub: structured rules + file extraction + optional reference PDF."""
    st.markdown(_policy_hub_nav_css(), unsafe_allow_html=True)
    st.markdown(_policy_hub_header_html(), unsafe_allow_html=True)

    _seg_key = f"{key_prefix}_policy_seg"
    view = st.segmented_control(
        "Policy sections",
        options=["structured", "file", "ref"],
        format_func=lambda k: {
            "structured": "Structured rules (by team)",
            "file": "Extract from uploaded file",
            "ref": "Reference document",
        }[k],
        default="structured",
        key=_seg_key,
        width="stretch",
        help="Switch between structured JSON-backed rules, optional file text extraction, and the on-disk PDF.",
    )
    active = view if view is not None else "structured"

    st.markdown("<div style='height:8px;'></div>", unsafe_allow_html=True)

    if active == "structured":
        st.markdown(
            _team_wrap("SMB", "SMB TEAM", "SMB Team — reps & manager", build_smb_inner_html()),
            unsafe_allow_html=True,
        )
        st.markdown(
            _team_wrap("AM", "AM TEAM", "AM Team", build_am_inner_html()),
            unsafe_allow_html=True,
        )
        st.markdown(
            _team_wrap("ENT", "ENT TEAM", "ENT Team — enterprise tiers", build_ent_inner_html()),
            unsafe_allow_html=True,
        )
        st.markdown(
            _team_wrap("OB", "OUTBOUND TEAM", "Outbound Team — meeting incentives", build_outbound_inner_html()),
            unsafe_allow_html=True,
        )
        st.markdown(
            _team_wrap("MN", "MONTHLY PLAN", "Monthly Plan — SMB & Account Management (April 2026)", build_monthly_plan_inner_html()),
            unsafe_allow_html=True,
        )

    elif active == "file":
        st.subheader("Upload PDF or PowerPoint")
        st.caption(
            "This does **not** replace policy JSON. Text is extracted for reading only; incentive math always follows "
            "`policy/commission_policy.json`."
        )
        st.markdown("<div style='height:10px;'></div>", unsafe_allow_html=True)
        up = st.file_uploader(
            "Commission deck (PDF or PPTX)",
            type=["pdf", "pptx"],
            key=f"{key_prefix}_policy_upload",
            help="Slides with heavy graphics may extract poorly — use Structured rules for authoritative tiers.",
        )
        if up is not None:
            raw = up.getvalue()
            text, err = extract_policy_document_text(raw, up.name)
            if err:
                st.error(err)
            elif text:
                st.success(f"Extracted **{len(text):,}** characters from **{_esc(up.name)}**.")
                st.markdown("<div style='height:14px;'></div>", unsafe_allow_html=True)
                buckets = bucket_extracted_lines(text)
                st.markdown("##### Routed excerpts (keyword heuristic)")
                st.caption("Spacing is widened between team columns for easier reading.")
                st.markdown("<div style='height:12px;'></div>", unsafe_allow_html=True)
                c1, gap, c2 = st.columns([1, 0.22, 1])
                with c1:
                    render_extracted_bucket_column("SMB", "SMB Team", buckets["SMB"])
                    st.markdown("<div style='height:16px;'></div>", unsafe_allow_html=True)
                    render_extracted_bucket_column("ENT", "ENT Team", buckets["ENT"])
                with gap:
                    st.markdown(
                        "<div style='min-height:1px;background:linear-gradient(180deg,transparent,#e2e8f0,transparent);"
                        "margin:8px 0;border-radius:2px;'></div>",
                        unsafe_allow_html=True,
                    )
                with c2:
                    render_extracted_bucket_column("AM", "AM Team", buckets["AM"])
                    st.markdown("<div style='height:16px;'></div>", unsafe_allow_html=True)
                    render_extracted_bucket_column("OB", "Outbound Team", buckets["OB"])
                other = buckets["OTHER"]
                if other:
                    with st.expander(f"Other lines ({len(other)} lines)", expanded=False):
                        st.text("\n".join(other[:400]) + ("\n…" if len(other) > 400 else ""))
                with st.expander("Full extracted text", expanded=False):
                    st.text(text[:120_000] + ("…\n[truncated]" if len(text) > 120_000 else ""))
            else:
                st.warning("No text could be extracted.")

    else:
        st.caption("Optional: original PDF on disk or from env — for download or quick open, not the primary policy view.")
        st.markdown("<div style='height:10px;'></div>", unsafe_allow_html=True)
        if pdf_path and pdf_path.is_file():
            st.info(f"File: `{pdf_path}`")
            try:
                data = pdf_path.read_bytes()
            except OSError as e:
                st.error(str(e))
            else:
                st.download_button(
                    label=f"Download {_esc(pdf_path.name)}",
                    data=data,
                    file_name=pdf_path.name,
                    mime="application/pdf",
                    key=f"{key_prefix}_dl_policy_pdf",
                )
                max_embed = 6 * 1024 * 1024
                if len(data) <= max_embed:
                    b64 = base64.b64encode(data).decode("utf-8")
                    st.markdown(
                        f'<iframe src="data:application/pdf;base64,{b64}" width="100%" height="780" type="application/pdf"></iframe>',
                        unsafe_allow_html=True,
                    )
                else:
                    st.info("PDF is large; use **Download** to open locally.")
        else:
            st.warning(
                "No reference PDF configured. Place a file under `policy/` or set **COMMISSION_POLICY_PDF** in `.env`."
            )
